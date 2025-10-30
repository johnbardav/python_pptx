"""
Contains all logic for generating application slides,
adapted from the reference scripts.
"""

import pandas as pd
import re
import os
import sys
from io import BytesIO
from unidecode import unidecode
from thefuzz import process, fuzz
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
import warnings
from sqlalchemy import text

# --- IMPORTAMOS LA PLANTILLA BASE ---
from .base_slide import create_base_slide

# Ocultar warnings de openpyxl
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

# --- CONFIGURACIÓN ---
OUTPUT_FOLDER = "outputs"
ICONS_FOLDER = "icons"
APP_COLUMN_NAME = "aplicacion_sistema"

# Constantes de BD (copiadas de db_loader.py)
TABLE_BUYER_BANK = "aplicaciones_buyer_bank"
TABLE_BOUGHT_BANK = "aplicaciones_bought_bank"

# Mapeo de Países a Iconos
COUNTRY_ICONS = {
    "Colombia (CO)": "co.png",
    "Panamá (PA)": "pa.png",
    "Costa Rica (CR)": "cr.png",
    "Honduras (HN)": "hn.png",
    "El Salvador (SV)": "sv.png",
    "EEUU - Miami (US)": "us.png",
}

# --- NUEVO: Orden de agrupación ---
COUNTRY_SORT_ORDER = {
    "Colombia (CO)": 1,
    "Panamá (PA)": 2,
    "Costa Rica (CR)": 3,
    "Honduras (HN)": 4,
    "El Salvador (SV)": 5,
    "EEUU - Miami (US)": 6,
}
BANK_SORT_ORDER = {"BuyerBank": 1, "BoughtBank": 2}
# --- Fin de Orden ---

# Parámetros (tomados de 1_generador_madurez_y_reportes.py)
SIMILARITY_THRESHOLD = 0.90
TECH_TRUNCATE_LENGTH = 33
ROW_HEIGHT = Cm(0.62)
TEXTBOX_HEIGHT = Cm(0.48)
ICON_SIZE = Cm(0.46)  # Tamaño estándar para iconos de criterios
FONT_SIZE = 8

INDICATOR_ICONS = {"si": "si.svg", "no": "no.svg", "parcial": "na.svg", "na": "na.svg"}
HEADER_LABELS = {
    "pais_icon": "País",
    "banco_icon": "",
    "aplicaciones": "Aplicaciones",
    "sas": "",
    "cloud": "",
    "cots": "",
    "regional": "",
    "tecnologia_subyacente": "Tecnología subyacente",
    "obsolescencia": "Obsolescencia",
    "escalabilidad": "Escalabilidad",
    "acople": "Acople",
    "estabilidad": "Estabilidad",
    "extensibilidad": "Extensibilidad",
    "seguridad": "Seguridad",
}
# Asegurarnos de usar los nombres de columna de la BD (con '_')
COLUMN_ORDER = [
    "pais_icon",
    "banco_icon",
    "aplicaciones",
    "sas",
    "cloud",
    "cots",
    "regional",
    "tecnologia_subyacente",  # <-- Movido
    "obsolescencia",
    "escalabilidad",
    "acople",
    "estabilidad",
    "extensibilidad",
    "seguridad",
]
COLUMN_WIDTHS = {
    "pais_icon": Cm(0.71),
    "banco_icon": Cm(0.55),
    "aplicaciones": Cm(6.70),
    "sas": Cm(0.6),
    "cloud": Cm(0.6),
    "cots": Cm(0.6),
    "regional": Cm(0.6),
    "tecnologia_subyacente": Cm(6.70),  # <-- Ancho igualado
    "obsolescencia": Cm(2.0),
    "escalabilidad": Cm(2.0),
    "acople": Cm(2.0),
    "estabilidad": Cm(2.0),
    "extensibilidad": Cm(2.0),
    "seguridad": Cm(2.0),
}

# Gaps entre las primeras columnas
GAPS = {
    "pais_icon": Cm(0.71),
    "banco_icon": Cm(0.42),
    # "aplicaciones" ya no se usa para calcular la siguiente posición
    "regional": Cm(0.05),  # Gap después de regional, antes de tech
}


# --- FUNCIONES DE LÓGICA (Adaptadas de la referencia) ---


def normalize_string(text):
    """(Función de masters/excel_loader.py)"""
    if not isinstance(text, str):
        return ""
    subscript_map = str.maketrans("₀₁₂₃₄₅₆₇₈₉", " " * 10)
    clean_text = text.translate(subscript_map)
    clean_text = unidecode(clean_text).lower()
    clean_text = re.sub(r"\s*\([^)]*\)\s*", " ", clean_text)
    noise_words = ["incluida en venta", "tsa", "no tsa"]
    for word in noise_words:
        clean_text = clean_text.replace(word, "")
    clean_text = re.sub(r"[^a-z0-9\s-]", "", clean_text)
    clean_text = clean_text.strip()
    return re.sub(r"\s+", " ", clean_text)


def get_value_from_row(row, col_name):
    """
    Obtiene un valor único de una fila (un pd.Series),
    manejando columnas duplicadas (ej. 'banco_1').
    """
    if not col_name or col_name not in row.index:
        return None

    value = row[col_name]

    # Si hay duplicados ('col', 'col_1'), row[col_name] puede devolver una Serie
    if isinstance(value, pd.Series):
        value = value.iloc[0]  # Tomar el primer valor

    if pd.notna(value):
        return str(value).strip()
    return None


def evaluar_criterios(row, bank_name, criteria_map):
    """
    Evalúa los criterios para una aplicación (fila) usando el mapeo
    proporcionado por 'config.py'.
    """
    resultados = {}

    # --- 1. OBSOLESCENCIA ---
    # Regla: 0 o vacío -> Vacío. "No obsoleto" -> Cumple. "Obsoleto" -> No Cumple.
    col_obs = criteria_map.get("obsolescencia")
    valor_obs_raw = get_value_from_row(row, col_obs)

    if not valor_obs_raw or valor_obs_raw == "0":
        resultados["obsolescencia"] = ""
    else:
        valor_lower = valor_obs_raw.lower()
        if "no obsoleto" in valor_lower:
            resultados["obsolescencia"] = "Cumple"
        elif "obsoleto" in valor_lower:
            resultados["obsolescencia"] = "No Cumple"
        else:
            resultados["obsolescencia"] = ""

    # --- 2. ESCALABILIDAD ---
    # Regla: "Si" -> Cumple. "No" -> No Cumple.
    col_esc = criteria_map.get("escalabilidad")
    valor_esc_raw = get_value_from_row(row, col_esc)

    if valor_esc_raw:
        valor_upper = valor_esc_raw.upper()
        if valor_upper == "SI":
            resultados["escalabilidad"] = "Cumple"
        elif valor_upper == "NO":
            resultados["escalabilidad"] = "No Cumple"
        else:
            resultados["escalabilidad"] = ""
    else:
        resultados["escalabilidad"] = ""

    # --- 3. ACOPLE ---
    # Regla: Siempre "Parcialmente" por defecto.
    resultados["acople"] = "Parcialmente"

    # --- 4. ESTABILIDAD ---
    # Regla: 0 or vacío -> Vacío. "Si" -> No Cumple. "No" -> Cumple.
    col_estab = criteria_map.get("estabilidad")
    valor_estab_raw = get_value_from_row(row, col_estab)

    if not valor_estab_raw or valor_estab_raw == "0":
        resultados["estabilidad"] = ""
    else:
        valor_upper = valor_estab_raw.upper()
        if valor_upper == "SI":
            resultados["estabilidad"] = "No Cumple"
        elif valor_upper == "NO":
            resultados["estabilidad"] = "Cumple"
        else:
            resultados["estabilidad"] = ""

    # --- 5. EXTENSIBILIDAD ---
    # Regla: "Cumple" -> Cumple. "Parcialmente" -> Parcialmente.
    col_ext = criteria_map.get("extensibilidad")
    valor_ext_raw = get_value_from_row(row, col_ext)

    if valor_ext_raw:
        valor_title = valor_ext_raw.title()
        if valor_title == "Cumple":
            resultados["extensibilidad"] = "Cumple"
        elif valor_title == "Parcialmente":
            resultados["extensibilidad"] = "Parcialmente"
        else:
            resultados["extensibilidad"] = ""
    else:
        resultados["extensibilidad"] = ""

    # --- 6. SEGURIDAD ---
    # Regla: 0 o vacío -> "" (Vacío). 1 o 2 -> No Cumple. 3 -> Parcialmente. 4 o 5 -> Cumple.
    col_seg = criteria_map.get("seguridad")
    valor_seg_raw = get_value_from_row(row, col_seg)

    try:
        if not valor_seg_raw or str(valor_seg_raw) == "0":
            resultados["seguridad"] = ""
        else:
            valor_seg_num = float(valor_seg_raw)
            if valor_seg_num in (1, 2):
                resultados["seguridad"] = "No Cumple"
            elif valor_seg_num == 3:
                resultados["seguridad"] = "Parcialmente"
            elif valor_seg_num in (4, 5):
                resultados["seguridad"] = "Cumple"
            else:
                resultados["seguridad"] = "N/A"
    except (ValueError, TypeError):
        # Si el valor es texto (ej. "N/A" del excel)
        if valor_seg_raw and valor_seg_raw.upper() == "N/A":
            resultados["seguridad"] = "N/A"
        else:
            resultados["seguridad"] = ""

    return resultados


def find_best_match(app_name, choices_dict):
    """(Función de 1_generador_madurez_y_reportes.py)"""
    normalized_name = normalize_string(app_name)
    if not choices_dict:
        return None
    if normalized_name in choices_dict:
        return choices_dict[normalized_name]

    # Usar 'process' de thefuzz para encontrar la mejor coincidencia
    match = process.extractOne(
        normalized_name, choices_dict.keys(), scorer=fuzz.token_set_ratio
    )

    if match and match[1] >= SIMILARITY_THRESHOLD * 100:
        return choices_dict[match[0]]

    return None


# --- FUNCIONES DE POWERPOINT (Adaptadas de la referencia) ---


def add_image(slide, img_path, x, y, height=ICON_SIZE, width=None):
    """
    (Función de 1_generador_madurez_y_reportes.py)
    Modificada para aceptar un ancho (width) opcional y asegurar alta calidad para SVGs.
    """
    if not os.path.exists(img_path):
        return None

    if width is None:
        width = height  # Si no se da ancho, asume un icono cuadrado

    try:
        if img_path.lower().endswith(".svg"):
            drawing = svg2rlg(img_path)
            if not drawing:
                return None
            png_buffer = BytesIO()
            # Asegurar alta DPI para la conversión SVG a PNG
            renderPM.drawToFile(drawing, png_buffer, fmt="PNG", dpi=600)
            png_buffer.seek(0)
            return slide.shapes.add_picture(
                png_buffer, x, y, width=width, height=height
            )
        elif img_path.lower().endswith((".png", ".jpg", ".jpeg")):
            # Para PNG/JPG existentes, los insertamos directamente.
            # La calidad de estos depende del archivo original.
            return slide.shapes.add_picture(img_path, x, y, width=width, height=height)
    except Exception:
        pass
    return None


def calculate_positions(start_x):
    """
    Calcula las posiciones X de inicio para cada columna.
    El lado izquierdo se calcula L-R y el lado derecho (criterios) R-L.
    """
    positions = {}

    # --- Lado Izquierdo (L-R) - Parte 1 ---
    cols_left_part1 = [
        "pais_icon",
        "banco_icon",
        "aplicaciones",
    ]

    current_x = start_x
    for col_name in cols_left_part1:
        positions[col_name] = current_x
        current_x += COLUMN_WIDTHS[col_name]
        if col_name in GAPS:
            current_x += GAPS[col_name]

    # --- Bloque de Iconos (Posición Absoluta) ---
    current_x = Cm(10.70)  # <-- POSICIÓN ACTUALIZADA

    positions["sas"] = current_x
    current_x += COLUMN_WIDTHS["sas"]

    positions["cloud"] = current_x
    current_x += COLUMN_WIDTHS["cloud"]

    positions["cots"] = current_x
    current_x += COLUMN_WIDTHS["cots"]

    positions["regional"] = current_x
    current_x += COLUMN_WIDTHS["regional"]

    # --- Lado Izquierdo (L-R) - Parte 2 (Tecnología) ---
    if "regional" in GAPS:
        current_x += GAPS["regional"]

    positions["tecnologia_subyacente"] = current_x

    # --- Lado Derecho (R-L) ---
    END_X = Cm(1.54) + Cm(30.8)

    cols_right_to_left = [
        "seguridad",
        "extensibilidad",
        "estabilidad",
        "acople",
        "escalabilidad",
        "obsolescencia",
    ]

    current_x_rtl = END_X
    for col_name in cols_right_to_left:
        current_x_rtl -= COLUMN_WIDTHS[col_name]
        positions[col_name] = current_x_rtl

    return positions


def draw_main_header(slide, x_positions, y_pos):
    """(Función de 1_generador_madurez_y_reportes.py)"""
    for col_name in COLUMN_ORDER:
        x, width = x_positions[col_name], COLUMN_WIDTHS[col_name]
        textbox = slide.shapes.add_textbox(x, y_pos, width, ROW_HEIGHT)
        tf = textbox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]

        # --- Lógica de alineación de cabecera ---
        if col_name in ["pais_icon", "aplicaciones", "tecnologia_subyacente"]:
            p.alignment = PP_ALIGN.LEFT
        else:
            p.alignment = PP_ALIGN.CENTER
        # --- Fin de la lógica de alineación ---

        run = p.add_run()
        run.text = HEADER_LABELS[col_name]
        run.font.bold = True
        run.font.size = Pt(FONT_SIZE)


# --- FUNCIÓN PRINCIPAL DEL MÓDULO (MODIFICADA) ---


def _get_sort_key(line_tuple):
    """Función helper para obtener la clave de ordenamiento País -> Banco"""
    # line_tuple es ((country, bank, app_name), original_line)
    country, bank, app_name = line_tuple[0]

    # Determinar bank_sort_name
    bank_sort_name = "BuyerBank" if "BUYERBANK" in bank.upper() else "BoughtBank"

    country_key = COUNTRY_SORT_ORDER.get(country, 99)  # 99 para desconocidos
    bank_key = BANK_SORT_ORDER.get(bank_sort_name, 99)

    return (
        country_key,
        bank_key,
        app_name.lower(),
    )  # Ordena alfabéticamente por nombre como 3er criterio


def generate_slide_for_subdomain(
    prs,
    subdomain_title,
    app_lines_data,
    df_buyer,
    choices_buyer,
    df_bought,
    choices_bought,
    criteria_map,
    db_engine,
):
    """
    Crea UN slide para un subdominio específico, usando la plantilla base.
    """

    print(f"    ▶️  Generando slide para el subdominio: {subdomain_title}...")

    # 1. CREAR EL SLIDE USANDO LA PLANTILLA BASE
    slide = create_base_slide(prs, subdomain_title, "")

    pending_lines = []

    # 2. DEFINIR COORDENADAS DE INICIO PARA EL CONTENIDO
    START_X = Cm(1.54)
    START_Y = Cm(5.22)  # Top del área de contenido

    # 3. DIBUJAR LA TABLA
    x_positions = calculate_positions(START_X)

    draw_main_header(slide, x_positions, START_Y)

    y_pos = START_Y + ROW_HEIGHT

    map_resultado_a_icono = {
        "Cumple": "si",
        "No Cumple": "no",
        "Parcialmente": "parcial",
        "N/A": "na",
    }

    # --- NUEVA LÓGICA DE ORDENAMIENTO ---
    sorted_app_lines_data = sorted(app_lines_data, key=_get_sort_key)
    current_country = None  # Para rastrear el grupo de país
    # --- FIN DE LÓGICA DE ORDENAMIENTO ---

    # Iterar sobre las líneas de app (AHORA ORDENADAS)
    for i, (line_data_tuple, original_line) in enumerate(sorted_app_lines_data):
        try:
            country, bank, app_name = line_data_tuple
            bank_upper = bank.upper()

            if "BUYERBANK" in bank_upper:
                target_df, target_choices = df_buyer, choices_buyer
                bank_name_eval = "BuyerBank"
                table_to_update = TABLE_BUYER_BANK
            elif "BOUGHTBANK" in bank_upper:
                target_df, target_choices = df_bought, choices_bought
                bank_name_eval = "BoughtBank"
                table_to_update = TABLE_BOUGHT_BANK
            else:
                print(
                    f"      ⚠️  Banco no reconocido (use 'BuyerBank' o 'BoughtBank'): '{original_line}'"
                )
                continue
        except Exception as e:
            print(f"      ❌ Error procesando línea '{original_line}': {e}")
            continue

        excel_match_name = find_best_match(app_name, target_choices)
        row = None

        if excel_match_name:
            # 1. ACTUALIZAR BASE DE DATOS
            try:
                update_stmt = text(
                    f"""
                    UPDATE {table_to_update}
                    SET mostrar_en_arquitectura_target = 'Si'
                    WHERE {APP_COLUMN_NAME} = :app_name
                    """
                )
                with db_engine.connect() as conn:
                    conn.execute(update_stmt, {"app_name": excel_match_name})
                    conn.commit()

            except Exception as e:
                print(f"      ⚠️  ERROR al actualizar BD para '{app_name}': {e}")

            # 2. OBTENER FILA PARA SLIDE
            row_df = (
                target_df[target_df[APP_COLUMN_NAME] == excel_match_name]
                if excel_match_name
                else pd.DataFrame()
            )
            row = row_df.iloc[0] if not row_df.empty else None

        else:
            # 2. APP NO ENCONTRADA, AGREGAR A PENDIENTES
            pending_lines.append(original_line)

        y_text_pos = y_pos + (ROW_HEIGHT - TEXTBOX_HEIGHT) / 2
        y_icon_pos_std = (
            y_pos + (ROW_HEIGHT - ICON_SIZE) / 2
        )  # Y-pos para iconos estándar

        # --- AÑADIR ICONO DE PAÍS (CONDICIONAL) ---
        if country != current_country:
            country_icon_name = COUNTRY_ICONS.get(country)
            if country_icon_name:
                icon_path = os.path.join(ICONS_FOLDER, country_icon_name)
                COUNTRY_ICON_H = Cm(0.51)
                COUNTRY_ICON_W = Cm(0.71)
                y_icon_pais = y_pos + (ROW_HEIGHT - COUNTRY_ICON_H) / 2
                x_icon_pais = x_positions["pais_icon"]
                add_image(
                    slide,
                    icon_path,
                    x_icon_pais,
                    y_icon_pais,
                    height=COUNTRY_ICON_H,
                    width=COUNTRY_ICON_W,
                )
            current_country = country  # Actualizar el país actual

        # --- AÑADIR ICONO DE BANCO (TAMAÑO PERSONALIZADO) ---
        if bank_name_eval == "BuyerBank":
            icon_path = os.path.join(ICONS_FOLDER, "dav.png")
        else:  # bank_name_eval == "BoughtBank"
            icon_path = os.path.join(ICONS_FOLDER, "sco.png")

        BANK_ICON_H = Cm(0.49)
        BANK_ICON_W = Cm(0.55)
        y_icon_banco = y_pos + (ROW_HEIGHT - BANK_ICON_H) / 2
        x_icon_banco = x_positions["banco_icon"] + (
            (COLUMN_WIDTHS["banco_icon"] - BANK_ICON_W) / 2
        )
        add_image(
            slide,
            icon_path,
            x_icon_banco,
            y_icon_banco,
            height=BANK_ICON_H,
            width=BANK_ICON_W,
        )

        # --- AÑADIR SHAPE DE APLICACIÓN ---
        x_app = x_positions["aplicaciones"]
        w_app = COLUMN_WIDTHS["aplicaciones"]
        y_app_shape = y_pos + (ROW_HEIGHT - TEXTBOX_HEIGHT) / 2

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_app, y_app_shape, w_app, TEXTBOX_HEIGHT
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("E3151E")
        shape.line.fill.background()  # Sin borde
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.text = app_name
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(FONT_SIZE)

        if row is not None:
            # Evaluar criterios usando el MAPA
            resultados_evaluacion = evaluar_criterios(row, bank_name_eval, criteria_map)

            # --- Iconos (SAS, COTS, CLOUD, REGIONAL) ---
            val_sas = get_value_from_row(row, criteria_map.get("icon_sas"))
            if val_sas and normalize_string(val_sas) == "si":
                icon_path = os.path.join(ICONS_FOLDER, "sass.png")
                x_icon = x_positions["sas"] + (COLUMN_WIDTHS["sas"] - ICON_SIZE) / 2
                add_image(slide, icon_path, x_icon, y_icon_pos_std, height=ICON_SIZE)

            val_custom = get_value_from_row(row, criteria_map.get("icon_cots"))
            if val_custom:
                valor_customizacion = normalize_string(val_custom)
                if valor_customizacion in ["cots", "cots con observacion"]:
                    icon_path = os.path.join(ICONS_FOLDER, "cots.svg")
                    x_icon = (
                        x_positions["cots"] + (COLUMN_WIDTHS["cots"] - ICON_SIZE) / 2
                    )
                    add_image(
                        slide, icon_path, x_icon, y_icon_pos_std, height=ICON_SIZE
                    )

            val_nube = get_value_from_row(row, criteria_map.get("icon_cloud"))
            if val_nube and normalize_string(val_nube) == "nube":
                icon_path = os.path.join(ICONS_FOLDER, "cloud.svg")
                x_icon = x_positions["cloud"] + (COLUMN_WIDTHS["cloud"] - ICON_SIZE) / 2
                add_image(slide, icon_path, x_icon, y_icon_pos_std, height=ICON_SIZE)

            val_bns_raw = get_value_from_row(row, criteria_map.get("icon_regional"))
            if val_bns_raw:
                valor_bns_norm = normalize_string(val_bns_raw)
                if "regional" in valor_bns_norm or "global" in valor_bns_norm:
                    icon_path = os.path.join(ICONS_FOLDER, "regional.svg")
                    x_icon = (
                        x_positions["regional"]
                        + (COLUMN_WIDTHS["regional"] - ICON_SIZE) / 2
                    )
                    add_image(
                        slide, icon_path, x_icon, y_icon_pos_std, height=ICON_SIZE
                    )

            # --- Tecnología ---
            tech_text = get_value_from_row(row, criteria_map.get("tecnologia"))
            if tech_text:
                # Reemplazar saltos de línea por espacios
                tech_text = " ".join(tech_text.split())

                if len(tech_text) > TECH_TRUNCATE_LENGTH:
                    tech_text = tech_text[:TECH_TRUNCATE_LENGTH] + "..."
                textbox_tec = slide.shapes.add_textbox(
                    x_positions["tecnologia_subyacente"],
                    y_text_pos,
                    COLUMN_WIDTHS["tecnologia_subyacente"],
                    TEXTBOX_HEIGHT,
                )
                tf_tec = textbox_tec.text_frame
                (
                    tf_tec.margin_left,
                    tf_tec.margin_right,
                    tf_tec.margin_top,
                    tf_tec.margin_bottom,
                ) = 0, 0, 0, 0
                tf_tec.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_tec = tf_tec.paragraphs[0]
                p_tec.text = tech_text
                p_tec.font.size = Pt(FONT_SIZE)

            # --- Iconos de Evaluación (Criterios) ---
            indicator_cols = [
                "obsolescencia",
                "escalabilidad",
                "acople",
                "estabilidad",
                "extensibilidad",
                "seguridad",
            ]

            for col_name in indicator_cols:
                resultado_evaluado = resultados_evaluacion.get(col_name)
                icono_key = map_resultado_a_icono.get(resultado_evaluado)
                if icono_key and icono_key in INDICATOR_ICONS:
                    icon_path = os.path.join(ICONS_FOLDER, INDICATOR_ICONS[icono_key])
                    x_icon = (
                        x_positions[col_name]
                        + (COLUMN_WIDTHS[col_name] - ICON_SIZE) / 2
                    )
                    add_image(
                        slide, icon_path, x_icon, y_icon_pos_std, height=ICON_SIZE
                    )

        # --- Añadir línea separadora (CONDICIONAL) ---
        line_y_pos = y_pos + ROW_HEIGHT

        # Lógica de Look-ahead
        is_last_item = i == len(sorted_app_lines_data) - 1
        next_country = None
        if not is_last_item:
            next_country, _, _ = sorted_app_lines_data[i + 1][
                0
            ]  # País del siguiente item

        line_width = Pt(0.5)  # Grosor estándar para todas

        if is_last_item or next_country != current_country:
            # Fin de grupo de país: Línea ANCHA
            line_x_start = x_positions["pais_icon"]  # Inicia en el icono de país
            line_x_end = line_x_start + Cm(30.81)
        else:
            # Registro normal: Línea CORTA
            line_x_start = x_positions["banco_icon"]  # Inicia en el icono de banco
            line_x_end = line_x_start + Cm(29.38)

        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, line_x_start, line_y_pos, line_x_end, line_y_pos
        )
        line = shape.line
        line.color.rgb = RGBColor.from_string("959092")
        line.width = line_width
        shape.shadow.inherit = False
        # --- Fin de línea separadora ---

        y_pos += ROW_HEIGHT

    return pending_lines
