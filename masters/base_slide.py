"""
Defines the base slide layouts for the Bank App Analysis project.
"""

import sys

try:
    from pptx.slide import Slide
    from pptx.presentation import Presentation
    from pptx.util import Cm
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: 'python-pptx' library not found. Please install it.")
    sys.exit(1)


def create_base_slide(prs: Presentation, title_text: str, content_text: str) -> Slide:
    """
    Adds a new slide to the presentation 'prs' using the standard
    base layout (title, toolbar, content area).

    Args:
        prs: The presentation object to add the slide to.
        title_text: The text for the title box.
        content_text: The text for the main content box.

    Returns:
        The slide object that was created.
    """

    # Use the "Blank" layout (index 6)
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        print("Warning: Blank layout (6) not found. Using layout 0 as fallback.")
        blank_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(blank_layout)

    # 1. Title configuration
    title_box = slide.shapes.add_textbox(Cm(1.54), Cm(0.48), Cm(30.78), Cm(2.03))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2. "Toolbar" shape
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.54), Cm(2.99), Cm(30.78), Cm(0.41))

    # 3. Content section
    content_box = slide.shapes.add_textbox(Cm(1.54), Cm(5.22), Cm(30.8), Cm(12.34))
    content_frame = content_box.text_frame
    content_frame.text = content_text

    return slide
